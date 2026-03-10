from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

TEMPLATE = 'scripts/template.pptx'
OUTPUT = 'deliverables/java-junit-training-deck.pptx'
PDF_OUT = 'deliverables/java-junit-training-deck.pdf'

prs = Presentation(TEMPLATE)

BLUE = RGBColor(26, 71, 140)
DEEP = RGBColor(20, 46, 88)
LIGHT = RGBColor(237, 242, 249)
MID = RGBColor(91, 123, 170)
ORANGE = RGBColor(230, 108, 10)
GOLD = RGBColor(191, 144, 0)
GRAY = RGBColor(96, 96, 96)
DARK = RGBColor(35, 35, 35)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(112, 173, 71)

TODAY = datetime.now().strftime('%Y/%m/%d')


def remove_slide(prs, index):
    slide_id = prs.slides._sldIdLst[index]
    rId = slide_id.rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[index]


def clear_text(shape):
    if not hasattr(shape, 'text_frame'):
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.clear()


def set_text(shape, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font='Microsoft YaHei'):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font_obj = run.font
    font_obj.name = font
    font_obj.size = Pt(size)
    font_obj.bold = bold
    font_obj.color.rgb = color
    p.alignment = align
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_box(slide, x, y, w, h, title, body, fill=LIGHT, line=MID, title_color=DEEP, body_color=DARK):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    shape.adjustments[0] = 0.08
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(6)

    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = 'Microsoft YaHei'
    r1.font.size = Pt(18)
    r1.font.bold = True
    r1.font.color.rgb = title_color

    p2 = tf.add_paragraph()
    p2.space_before = Pt(5)
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = 'Microsoft YaHei'
    r2.font.size = Pt(12.5)
    r2.font.color.rgb = body_color
    return shape


def add_bullets(slide, x, y, w, h, items, size=16, color=DARK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(7)
        for run in p.runs:
            run.font.name = 'Microsoft YaHei'
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return box


def add_code(slide, x, y, w, h, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DEEP
    shape.line.color.rgb = DEEP
    shape.adjustments[0] = 0.05
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = 'Consolas'
    run.font.size = Pt(12)
    run.font.color.rgb = WHITE
    return shape


def set_title(slide, title):
    shape = slide.shapes.title
    if shape is not None:
        set_text(shape, title, size=24, bold=True, color=WHITE)


def add_note_bar(slide, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.05), Inches(12.0), Inches(0.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.color.rgb = BLUE
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = 'Microsoft YaHei'
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE


# Keep original cover slide only, remove other template sample slides.
for idx in range(len(prs.slides) - 1, 0, -1):
    remove_slide(prs, idx)

# Slide 1 cover based on original slide.
cover = prs.slides[0]
set_text(cover.shapes[2], 'Java JUnit 单元测试培训', size=26, bold=True, color=DARK, align=PP_ALIGN.CENTER)
set_text(cover.shapes[1], '汇 报 人：晏斐\n汇报时间：' + TODAY, size=11, bold=False, color=DARK, align=PP_ALIGN.LEFT)
if hasattr(cover.shapes[3], 'text_frame'):
    clear_text(cover.shapes[3])
if hasattr(cover.shapes[4], 'text_frame'):
    set_text(cover.shapes[4], 'Maven + JDK 1.8 + JUnit 5 + Mockito + Spring Boot + MyBatis', size=15, bold=False, color=GRAY, align=PP_ALIGN.CENTER)

# Slide 2 objectives
slide = prs.slides.add_slide(prs.slide_layouts[2])
set_title(slide, '培训目标与学习收益')
add_box(slide, 0.8, 1.1, 3.8, 1.55, '为什么学', '降低回归风险，支撑重构，建立团队统一的测试语义。', fill=LIGHT, line=MID)
add_box(slide, 4.8, 1.1, 3.8, 1.55, '学什么', 'JUnit 5、Mockito、Spring Boot 测试分层、MyBatis 数据访问测试。', fill=LIGHT, line=MID)
add_box(slide, 8.8, 1.1, 3.8, 1.55, '学完能做什么', '能围绕核心业务写稳定测试，并将报告和覆盖率纳入工程流程。', fill=LIGHT, line=MID)
add_bullets(slide, 0.95, 3.15, 11.3, 2.0, [
    '从“会写 @Test”升级为“知道什么该测、怎么测、测到什么程度”。',
    '所有内容都能回到项目代码和测试类中直接演示。',
    '培训重点是持续写对测试，而不是临时补几个用例。'
], size=16)
add_note_bar(slide, '对应章节：01-培训目标.md')

# Slide 3 agenda
slide = prs.slides.add_slide(prs.slide_layouts[2])
set_title(slide, '目录')
sub = slide.shapes.add_textbox(Inches(0.85), Inches(1.0), Inches(2.2), Inches(0.3))
set_text(sub, 'CONTENTS', size=15, bold=True, color=MID)
add_box(slide, 0.82, 1.55, 11.4, 3.95, '培训结构', '1. 基础入门：环境、JUnit 核心写法、生命周期\n\n2. JUnit 进阶：参数化、Mock、高级技巧\n\n3. 框架实战：Spring Boot 分层、MyBatis CRUD 与集成\n\n4. 工程落地：最佳实践、报告与覆盖率、团队推进', fill=LIGHT, line=MID)

# Slide 4 env
slide = prs.slides.add_slide(prs.slide_layouts[2])
set_title(slide, '环境与项目结构')
add_box(slide, 0.8, 1.05, 3.9, 1.35, '基础环境', 'JDK 1.8\nMaven 构建\n标准 src/main / src/test 结构', fill=LIGHT, line=MID)
add_code(slide, 0.82, 2.7, 3.9, 2.2, 'mvn test\n\nmvn -Dtest=CalculatorTest test\n\nmvn clean verify')
add_box(slide, 5.1, 1.05, 3.45, 1.1, 'core', '纯 Java 业务逻辑，适合基础单测与参数化测试', fill=LIGHT, line=MID)
add_box(slide, 8.8, 1.05, 3.45, 1.1, 'service', '依赖协作、Mockito 与外部接口 Mock 示例', fill=LIGHT, line=MID)
add_box(slide, 5.1, 2.45, 3.45, 1.1, 'spring', 'Spring Boot Controller / Service / Mapper 测试', fill=LIGHT, line=MID)
add_box(slide, 8.8, 2.45, 3.45, 1.1, 'test/resources', 'SQL 初始化脚本与测试资源', fill=LIGHT, line=MID)
add_bullets(slide, 5.05, 4.05, 7.15, 1.4, [
    '目录结构清晰，IDE、Maven 和 CI 才能按约定工作。',
    'H2、Surefire、JaCoCo 让本地与流水线环境保持一致。'
], size=15)
add_note_bar(slide, '对应章节：02-环境与项目结构.md')

# Slide 5 junit core
slide = prs.slides.add_slide(prs.slide_layouts[2])
set_title(slide, 'JUnit 核心写法：AAA + 三类路径')
add_code(slide, 0.82, 1.1, 4.15, 2.4, '// Arrange\nCalculator calculator = new Calculator();\n\n// Act\nint result = calculator.add(3, 4);\n\n// Assert\nassertEquals(7, result);')
add_box(slide, 5.25, 1.1, 2.2, 1.2, '正常路径', '输入合法时，验证主业务逻辑。', fill=LIGHT, line=MID)
add_box(slide, 7.75, 1.1, 2.2, 1.2, '异常路径', '非法输入或状态不满足时，验证异常。', fill=LIGHT, line=MID)
add_box(slide, 10.25, 1.1, 2.0, 1.2, '边界路径', '0、空串、null、最大值与阈值附近。', fill=LIGHT, line=MID)
add_bullets(slide, 5.25, 2.8, 6.9, 1.8, [
    '常用断言：assertEquals、assertThrows、assertTimeout。',
    '培训讲解顺序建议：正常 -> 异常 -> 边界。',
    '对应示例类：CalculatorTest、TimeoutAndDisabledTest。'
], size=15)
add_note_bar(slide, '对应章节：03-JUnit核心写法.md')

# Slide 6 lifecycle + param basic
slide = prs.slides.add_slide(prs.slide_layouts[2])
set_title(slide, '生命周期与参数化测试基础')
add_box(slide, 0.82, 1.05, 3.8, 1.4, '生命周期', '@BeforeEach / @AfterEach 管理初始化与清理；@Nested 按业务场景分组。', fill=LIGHT, line=MID)
add_box(slide, 4.85, 1.05, 3.8, 1.4, '参数化价值', '把固定测试逻辑与变化测试数据分离，减少重复代码。', fill=LIGHT, line=MID)
add_box(slide, 8.88, 1.05, 3.35, 1.4, '常见参数源', '@CsvSource\n@ValueSource\n@NullSource', fill=LIGHT, line=MID)
add_code(slide, 0.82, 2.9, 5.3, 1.95, '@ParameterizedTest\n@CsvSource({" TOM ,Tom", "jErrY,Jerry"})\nvoid shouldNormalize(String raw, String expected) {\n    assertEquals(expected, normalizer.normalizeName(raw));\n}')
add_bullets(slide, 6.45, 2.95, 5.7, 1.8, [
    '测试必须彼此独立，不能依赖执行顺序。',
    '参数化适合“同一规则、多组数据”，不适合完全不同的测试意图。'
], size=15)
add_note_bar(slide, '对应章节：04-生命周期与测试组织.md，05-参数化测试基础.md')

# Slide 7 advanced param + advanced junit
slide = prs.slides.add_slide(prs.slide_layouts[2])
set_title(slide, '参数化测试进阶与高级技巧')
add_box(slide, 0.82, 1.05, 5.1, 1.3, '@MethodSource', '适合多参数与复杂业务规则；每组 Arguments 应完整表达一个业务事实。', fill=LIGHT, line=MID)
add_code(slide, 0.82, 2.55, 5.1, 2.15, 'static Stream<Arguments> validCases() {\n    return Stream.of(\n        arguments(100.0, "VIP", 80.0),\n        arguments(200.0, "NORMAL", 200.0)\n    );\n}')
add_box(slide, 6.2, 1.05, 2.0, 1.05, 'assertAll', '聚合同一业务动作下的多个断言', fill=LIGHT, line=MID)
add_box(slide, 8.45, 1.05, 2.0, 1.05, 'assertDoesNotThrow', '显式声明合法输入不应抛错', fill=LIGHT, line=MID)
add_box(slide, 10.7, 1.05, 1.6, 1.05, 'assumeTrue', '条件不满足时跳过', fill=LIGHT, line=MID)
add_box(slide, 6.2, 2.4, 6.1, 1.15, '@RepeatedTest / @TestFactory', '增强表达力，但不应为了“更高级”而使用。', fill=LIGHT, line=MID)
add_bullets(slide, 6.2, 3.9, 6.1, 1.4, [
    '@TestFactory 更适合轻量、纯逻辑动态用例。',
    '固定数据、明确生命周期控制的场景优先 @ParameterizedTest。'
], size=14.5)
add_note_bar(slide, '对应章节：06-参数化测试进阶.md，08-高级技巧.md')

# Slide 8 mock basics
slide = prs.slides.add_slide(prs.slide_layouts[2])
set_title(slide, 'Mock 与依赖隔离')
add_box(slide, 0.82, 1.05, 2.5, 1.15, '@Mock', '替换仓储、消息、邮件、HTTP 等外部依赖', fill=LIGHT, line=MID)
add_box(slide, 3.55, 1.05, 2.6, 1.15, '@InjectMocks', '把依赖注入被测服务，聚焦当前类职责', fill=LIGHT, line=MID)
add_box(slide, 6.38, 1.05, 2.7, 1.15, 'when / then', '控制返回、异常或动态响应', fill=LIGHT, line=MID)
add_box(slide, 9.32, 1.05, 2.9, 1.15, 'verify', '验证结果之外的副作用与交互', fill=LIGHT, line=MID)
add_code(slide, 0.82, 2.55, 5.5, 1.95, 'when(userRepository.findByEmail("new@company.com"))\n    .thenReturn(null);\nverify(emailSender).sendWelcomeEmail("new@company.com");')
add_bullets(slide, 6.62, 2.6, 5.45, 1.7, [
    'ArgumentCaptor：校验传入 save(...) 的对象字段。',
    'InOrder：校验调用顺序。',
    'verifyNoMoreInteractions：限制额外交互。'
], size=15)
add_note_bar(slide, '对应章节：07-Mock与依赖隔离.md（服务层基础示例）')

# Slide 9 external api mock
slide = prs.slides.add_slide(prs.slide_layouts[2])
set_title(slide, '外部接口 Mock：从单次调用到批量聚合')
add_box(slide, 0.82, 1.05, 2.7, 1.2, '基础场景', '正常返回、超时、服务不可用', fill=LIGHT, line=MID)
add_box(slide, 3.8, 1.05, 2.7, 1.2, '异常分类', '连接异常、业务错误响应、运行时异常', fill=LIGHT, line=MID)
add_box(slide, 6.78, 1.05, 2.7, 1.2, '容错', '空响应、未知状态、输入校验', fill=LIGHT, line=MID)
add_box(slide, 9.76, 1.05, 2.46, 1.2, '聚合', '批量查询下的部分成功部分失败', fill=LIGHT, line=MID)
add_code(slide, 0.82, 2.6, 5.7, 1.7, 'when(externalOrderClient.queryStatus("ORD-1006"))\n    .thenThrow(new SocketTimeoutException("Read timed out"))\n    .thenReturn(new OrderStatusResponse("ORD-1006", "PAID"));')
add_bullets(slide, 6.8, 2.65, 5.3, 1.9, [
    'thenThrow(...).thenReturn(...) 适合讲“瞬时失败重试”。',
    '批量查询要验证：单项失败不拖垮整批、超时项重试次数、非法输入不触发远端调用。',
    '对应示例：OrderStatusServiceExternalCallTest。'
], size=14.2)
add_note_bar(slide, '对应章节：07-Mock与依赖隔离.md（外部接口基础 + 进阶场景）')

# Slide 10 spring boot
slide = prs.slides.add_slide(prs.slide_layouts[2])
set_title(slide, 'Spring Boot 测试分层')
add_box(slide, 0.82, 1.15, 3.6, 1.45, '纯单元测试', '不启动 Spring，执行最快，适合规则判断、字符串处理和参数校验。', fill=LIGHT, line=MID)
add_box(slide, 4.85, 1.15, 3.6, 1.45, '切片测试', '@WebMvcTest / @MybatisTest，只验证某一层框架行为。', fill=LIGHT, line=MID)
add_box(slide, 8.88, 1.15, 3.35, 1.45, '全上下文测试', '@SpringBootTest，验证多层协作、自动装配和配置链路。', fill=LIGHT, line=MID)
add_bullets(slide, 0.95, 3.1, 5.6, 1.8, [
    '只关心业务逻辑：优先纯单元测试。',
    '只关心某一层框架行为：用切片测试。',
    '关心跨层协作：再用 @SpringBootTest。'
], size=15)
add_code(slide, 6.85, 3.0, 5.25, 1.8, '@WebMvcTest(GreetingController.class)\nclass GreetingControllerWebMvcTest { ... }\n\n@SpringBootTest\nclass EmployeeServiceSpringBootTest { ... }')
add_note_bar(slide, '对应章节：09-SpringBoot测试分层.md')

# Slide 11 mybatis
slide = prs.slides.add_slide(prs.slide_layouts[2])
set_title(slide, 'MyBatis：Mapper CRUD 与 Service 集成')
add_box(slide, 0.82, 1.05, 5.4, 1.25, 'Mapper 层关注点', 'insert、selectById、update、delete、selectAll 都要验证，写入操作建议“写后再读”。', fill=LIGHT, line=MID)
add_box(slide, 6.5, 1.05, 5.72, 1.25, '为什么用 H2 + @Sql', '保证本地、CI、培训环境一致，并隔离每个测试的数据状态。', fill=LIGHT, line=MID)
add_code(slide, 0.82, 2.65, 5.1, 1.55, '@MybatisTest\n@Sql(scripts = "/sql/employee-schema.sql",\n     executionPhase = Sql.ExecutionPhase.BEFORE_TEST_METHOD)')
add_bullets(slide, 6.45, 2.7, 5.6, 1.5, [
    'Mapper 测试看 SQL 与映射是否正确。',
    'Service 集成测试看 Spring 注入与完整业务链路是否成立。'
], size=15)
add_box(slide, 0.82, 4.6, 11.4, 1.05, '示例类', 'EmployeeMapperMybatisTest、EmployeeServiceSpringBootTest、employee-schema.sql', fill=LIGHT, line=MID)
add_note_bar(slide, '对应章节：10-MyBatisMapper层CRUD测试.md，11-MyBatisService集成测试.md')

# Slide 12 best practices
slide = prs.slides.add_slide(prs.slide_layouts[2])
set_title(slide, '最佳实践与测试报告')
add_box(slide, 0.82, 1.05, 5.45, 2.8, '最佳实践', '1. 一个测试只验证一个明确行为\n2. 命名体现业务语义\n3. 断言聚焦业务结果\n4. 依赖通过构造器注入\n5. 先设计可测试边界，再进入测试阶段', fill=LIGHT, line=MID)
add_box(slide, 6.58, 1.05, 5.64, 2.8, '报告与覆盖率', 'mvn clean verify\nSurefire 看 Failures / Errors / Skipped\nJaCoCo 看关键类和异常分支覆盖率\n覆盖率高不等于测试质量高', fill=LIGHT, line=MID)
add_bullets(slide, 0.95, 4.15, 11.2, 1.1, [
    '当前项目可以直接用 OrderStatusServiceExternalCallTest 观察异常分支覆盖率变化。',
    'CI 建议：固定执行 mvn clean verify，收集 Surefire XML 与 jacoco.xml。'
], size=15)
add_note_bar(slide, '对应章节：12-最佳实践与反模式.md，13-测试报告与覆盖率报告生成.md')

# Slide 13 rollout
slide = prs.slides.add_slide(prs.slide_layouts[2])
set_title(slide, '团队落地建议')
add_box(slide, 0.82, 1.15, 2.75, 1.3, '第一步', '统一最低标准：新增核心逻辑必须带测试。', fill=LIGHT, line=MID)
add_box(slide, 3.85, 1.15, 2.75, 1.3, '第二步', '把测试纳入 Code Review，看场景和断言质量。', fill=LIGHT, line=MID)
add_box(slide, 6.88, 1.15, 2.75, 1.3, '第三步', '接入 CI、报告和覆盖率阈值。', fill=LIGHT, line=MID)
add_box(slide, 9.91, 1.15, 2.35, 1.3, '第四步', '建立持续改进机制。', fill=LIGHT, line=MID)
add_bullets(slide, 0.95, 3.15, 11.1, 1.7, [
    '优先覆盖核心业务流程、高频变更模块、缺陷历史较多模块。',
    '老项目不要一开始追求全补齐，应采取“新增必测、修改即补”的增量策略。',
    '培训的终点不是听懂，而是形成可执行、可评审、可持续的测试习惯。'
], size=15)
add_note_bar(slide, '对应章节：14-QA与落地建议.md')

# Final thanks
slide = prs.slides.add_slide(prs.slide_layouts[2])
set_title(slide, '')
tb = slide.shapes.add_textbox(Inches(2.6), Inches(2.1), Inches(8.2), Inches(0.7))
set_text(tb, 'THANK YOU！', size=28, bold=True, color=DEEP, align=PP_ALIGN.CENTER, font='Times New Roman')
tb2 = slide.shapes.add_textbox(Inches(2.15), Inches(3.1), Inches(9.0), Inches(0.5))
set_text(tb2, '欢迎交流：基于示例工程继续扩展团队测试规范', size=17, bold=False, color=DEEP, align=PP_ALIGN.CENTER)
tb3 = slide.shapes.add_textbox(Inches(4.9), Inches(4.0), Inches(3.4), Inches(0.4))
set_text(tb3, '谢谢聆听', size=18, bold=False, color=GRAY, align=PP_ALIGN.CENTER)

prs.save(OUTPUT)
print(OUTPUT)
